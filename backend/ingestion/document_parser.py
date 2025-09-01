import os
import pdfplumber
from pptx import Presentation
from typing import List, Dict, Any
import re
import platform
from backend.config import settings

# Import for .ppt support (Windows only)
if platform.system() == "Windows":
    try:
        import win32com.client
        PPT_SUPPORT = True
    except ImportError:
        PPT_SUPPORT = False
        print("Warning: pywin32 not available. .ppt files will not be supported.")
else:
    PPT_SUPPORT = False

class DocumentParser:
    """Handles parsing of PDF and PowerPoint documents with overlapping chunking."""
    
    def __init__(self):
        self.chunk_size = 500  # words per chunk
        self.chunk_overlap = 50  # overlapping words
    
    def parse_pdf(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse PDF file and extract text with metadata."""
        chunks = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text:
                        # Clean and normalize text
                        text = self._clean_text(text)
                        
                        # Create chunks with overlapping
                        page_chunks = self._create_overlapping_chunks(
                            text, 
                            f"Page {page_num}",
                            file_path
                        )
                        chunks.extend(page_chunks)
        
        except Exception as e:
            print(f"Error parsing PDF {file_path}: {str(e)}")
        
        return chunks
    
    def parse_powerpoint(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse PowerPoint file and extract text with metadata."""
        chunks = []
        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_extension == '.ppt' and PPT_SUPPORT:
                # Handle .ppt files using win32com
                chunks = self._parse_ppt_file(file_path)
            else:
                # Handle .pptx files using python-pptx
                chunks = self._parse_pptx_file(file_path)
        
        except Exception as e:
            print(f"Error parsing PowerPoint {file_path}: {str(e)}")
        
        return chunks
    
    def _parse_pptx_file(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse .pptx file using python-pptx."""
        chunks = []
        
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                # Combine all text from the slide
                full_text = " ".join(slide_text)
                full_text = self._clean_text(full_text)
                
                # Create chunks with overlapping
                slide_chunks = self._create_overlapping_chunks(
                    full_text,
                    f"Slide {slide_num}",
                    file_path
                )
                chunks.extend(slide_chunks)
        
        return chunks
    
    def _parse_ppt_file(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse .ppt file using win32com (Windows only)."""
        chunks = []
        powerpoint = None
        presentation = None
        
        try:
            # Convert to absolute path
            abs_path = os.path.abspath(file_path)
            
            # Use PowerPoint application to open the file
            powerpoint = win32com.client.Dispatch("PowerPoint.Application")
            # Don't try to hide PowerPoint - it may not be allowed
            
            presentation = powerpoint.Presentations.Open(abs_path, ReadOnly=True)
            
            for slide_num in range(1, presentation.Slides.Count + 1):
                try:
                    slide = presentation.Slides(slide_num)
                    slide_text = []
                    
                    # Extract text from all shapes in the slide
                    for shape_idx in range(1, slide.Shapes.Count + 1):
                        try:
                            shape = slide.Shapes(shape_idx)
                            
                            # Check if shape has text
                            if hasattr(shape, "TextFrame") and shape.TextFrame:
                                try:
                                    text = shape.TextFrame.TextRange.Text.strip()
                                    if text:
                                        slide_text.append(text)
                                except:
                                    # Skip this text frame if there's an error
                                    continue
                        except:
                            # Skip this shape if there's an error
                            continue
                    
                    if slide_text:
                        # Combine all text from the slide
                        full_text = " ".join(slide_text)
                        full_text = self._clean_text(full_text)
                        
                        # Create chunks with overlapping
                        slide_chunks = self._create_overlapping_chunks(
                            full_text,
                            f"Slide {slide_num}",
                            file_path
                        )
                        chunks.extend(slide_chunks)
                
                except Exception as slide_error:
                    print(f"   Warning: Error processing slide {slide_num}: {slide_error}")
                    continue
        
        except Exception as e:
            print(f"   Error opening PowerPoint file: {e}")
        
        finally:
            # Clean up
            try:
                if presentation:
                    presentation.Close()
                if powerpoint:
                    powerpoint.Quit()
            except:
                pass
        
        return chunks
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove special characters but keep punctuation
        text = re.sub(r'[^\w\s\.\,\!\?\;\:\-\(\)\[\]]', '', text)
        
        # Normalize line breaks
        text = text.replace('\n', ' ').replace('\r', ' ')
        
        return text.strip()
    
    def _create_overlapping_chunks(self, text: str, section: str, source_file: str) -> List[Dict[str, Any]]:
        """Create overlapping text chunks."""
        chunks = []
        words = text.split()
        
        if len(words) <= self.chunk_size:
            # If text is smaller than chunk size, create single chunk
            chunks.append({
                'text': text,
                'section': section,
                'source_file': os.path.basename(source_file),
                'chunk_id': f"{section}_0",
                'word_count': len(words)
            })
        else:
            # Create overlapping chunks
            start = 0
            chunk_id = 0
            
            while start < len(words):
                end = min(start + self.chunk_size, len(words))
                chunk_words = words[start:end]
                chunk_text = ' '.join(chunk_words)
                
                chunks.append({
                    'text': chunk_text,
                    'section': section,
                    'source_file': os.path.basename(source_file),
                    'chunk_id': f"{section}_{chunk_id}",
                    'word_count': len(chunk_words)
                })
                
                # Move start position with overlap
                start += self.chunk_size - self.chunk_overlap
                chunk_id += 1
        
        return chunks
    
    def parse_document(self, file_path: str) -> List[Dict[str, Any]]:
        """Parse document based on file extension."""
        file_extension = os.path.splitext(file_path)[1].lower()
        
        if file_extension == '.pdf':
            return self.parse_pdf(file_path)
        elif file_extension == '.pptx':
            return self.parse_powerpoint(file_path)
        elif file_extension == '.ppt':
            if not PPT_SUPPORT:
                raise ValueError(
                    f".ppt files are not supported on this platform. "
                    f"Please convert your .ppt file to .pptx format before uploading."
                )
            return self.parse_powerpoint(file_path)
        else:
            raise ValueError(f"Unsupported file format: {file_extension}") 